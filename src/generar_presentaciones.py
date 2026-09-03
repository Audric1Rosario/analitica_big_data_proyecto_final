"""
Generador de Presentaciones PowerPoint (.pptx) Profesionales
Proyecto Final: Aplicaciones Analíticas de Big Data (UAPA)
Equipo: Audric Rosario & Orlando Benítez
Empresa: Claro República Dominicana
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Paleta de Colores Corporativa
COLOR_CLARO_RED = RGBColor(218, 41, 28)       # #DA291C
COLOR_DARK_NAVY = RGBColor(15, 23, 42)        # #0F172A
COLOR_CARD_BG   = RGBColor(30, 41, 59)        # #1E293B
COLOR_WHITE     = RGBColor(255, 255, 255)
COLOR_TEXT_MUTED= RGBColor(148, 163, 184)     # #94A3B8
COLOR_GREEN     = RGBColor(16, 185, 129)      # #10B981
COLOR_LIGHT_BG  = RGBColor(248, 250, 252)     # #F8FAFC
COLOR_BORDER    = RGBColor(51, 65, 85)        # #334155


def crear_diapositiva_base(prs, titulo: str, subtitulo: str = "", orador: str = "", apoyo_visual: str = ""):
    """Crea una diapositiva con estilo corporativo consistente en formato 16:9."""
    slide_layout = prs.slide_layouts[6]  # Diapositiva en blanco
    slide = prs.slides.add_slide(slide_layout)

    # Fondo general oscuro de alta gama
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = COLOR_DARK_NAVY
    background.line.color.rgb = COLOR_DARK_NAVY

    # Barra superior de acento Claro Red
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_CLARO_RED
    top_bar.line.color.rgb = COLOR_CLARO_RED

    # Título y Subtítulo
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(8.5), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.bold = True
    p.font.size = Pt(26)
    p.font.color.rgb = COLOR_WHITE

    if subtitulo:
        p2 = tf.add_paragraph()
        p2.text = subtitulo
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_TEXT_MUTED

    # Badge de Orador y Apoyo Visual (Esquina Superior Derecha)
    if orador or apoyo_visual:
        badge_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2), Inches(0.45), Inches(3.4), Inches(0.85))
        badge_box.fill.solid()
        badge_box.fill.fore_color.rgb = COLOR_CARD_BG
        badge_box.line.color.rgb = COLOR_BORDER
        btf = badge_box.text_frame
        btf.word_wrap = True
        bp1 = btf.paragraphs[0]
        bp1.text = f"🎤 {orador}" if orador else ""
        bp1.font.bold = True
        bp1.font.size = Pt(11)
        bp1.font.color.rgb = COLOR_CLARO_RED
        if apoyo_visual:
            bp2 = btf.add_paragraph()
            bp2.text = f"🎬 {apoyo_visual}"
            bp2.font.size = Pt(9.5)
            bp2.font.color.rgb = COLOR_TEXT_MUTED

    # Pie de página discreto
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.35))
    ftf = footer_box.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "UAPA • Aplicaciones Analíticas de Big Data • Proyecto Final: Claro República Dominicana"
    fp.font.size = Pt(9.5)
    fp.font.color.rgb = COLOR_TEXT_MUTED

    return slide


def agregar_tarjeta(slide, left, top, width, height, titulo, contenido, color_acento=COLOR_CLARO_RED, bg_color=COLOR_CARD_BG):
    """Crea un contenedor de tarjeta estilizado."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = COLOR_BORDER

    # Barra lateral de acento de la tarjeta
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top + 0.15), Inches(0.08), Inches(height - 0.3))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color_acento
    accent.line.color.rgb = color_acento

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.18)

    p1 = tf.paragraphs[0]
    p1.text = titulo
    p1.font.bold = True
    p1.font.size = Pt(14)
    p1.font.color.rgb = COLOR_WHITE

    for item in contenido:
        p = tf.add_paragraph()
        p.text = f"• {item}" if not item.startswith("  ") else item
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MUTED
        p.space_after = Pt(4)


# ==============================================================================
# 1. PRESENTACIÓN OFICIAL PARA EL VIDEO (8 MINUTOS - 8 DIAPOSITIVAS)
# ==============================================================================
def generar_presentacion_video(output_path: str = "PRESENTACION_VIDEO_OFICIAL.pptx"):
    print("[PPTX] Generando presentación para el Video Oficial (8 minutos)...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- SLIDE 1: PORTADA ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_DARK_NAVY
    bg1.line.color.rgb = COLOR_DARK_NAVY

    top1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.15))
    top1.fill.solid()
    top1.fill.fore_color.rgb = COLOR_CLARO_RED

    tx = slide1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.3), Inches(4.8))
    tf = tx.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "UNIVERSIDAD ABIERTA PARA ADULTOS (UAPA)\nMAESTRÍA EN ANALÍTICA DE BIG DATA E INTELIGENCIA DE NEGOCIOS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MUTED

    p2 = tf.add_paragraph()
    p2.text = "Auditoría de Experiencia del Cliente y Sentimiento de Marca en Telecomunicaciones vía YouTube Data API v3"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.space_before = Pt(18)

    p3 = tf.add_paragraph()
    p3.text = "Caso de Estudio: Claro República Dominicana (@clarord) — Procesamiento de Lenguaje Natural (NLP)"
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLOR_CLARO_RED
    p3.space_before = Pt(8)

    p4 = tf.add_paragraph()
    p4.text = "Facilitador: Luis Eduardo Bayonet Robles\nEquipo de Trabajo:\n• Audric André Rosario Rosario (100089140) — Lead Data Engineering & NLP Modeling\n• Orlando Benítez Ventura (100090873) — Lead Business Intelligence & Strategy"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_WHITE
    p4.space_before = Pt(20)

    slide1.notes_slide.notes_text_frame.text = (
        "[00:00 - 00:35] ORLANDO: Saludos cordiales al facilitador y a los compañeros. "
        "Mi nombre es Orlando Benítez y junto a Audric Rosario presentamos nuestro proyecto final..."
    )

    # --- SLIDE 2: EL PROBLEMA DE NEGOCIO ---
    slide2 = crear_diapositiva_base(
        prs, "El Reto Empresarial: Puntos Ciegos en la Voz del Cliente",
        "Por qué las encuestas tradicionales fallan en telecomunicaciones",
        orador="Orlando Benítez", apoyo_visual="Mostrar Diapositiva 2 completa"
    )
    agregar_tarjeta(slide2, 0.8, 1.8, 3.7, 4.8, "Encuestas Tradicionales", [
        "Tasa de respuesta < 4% (CSAT / SMS).",
        "Alta latencia: Reportes mensuales/trimestrales.",
        "Sesgo emocional: Solo responden los extremos.",
        "Costo elevado y muestreo reducido."
    ], COLOR_BORDER)
    agregar_tarjeta(slide2, 4.8, 1.8, 3.7, 4.8, "Voz en Redes (YouTube)", [
        "Comentarios espontáneos y no solicitados.",
        "Tiempo real durante averías masivas.",
        "Reseñas técnicas de fibra y velocidad 5G.",
        "Punto ciego: No se analizan sistemáticamente."
    ], COLOR_CLARO_RED)
    agregar_tarjeta(slide2, 8.8, 1.8, 3.7, 4.8, "Impacto Comercial (Churn)", [
        "El costo de adquirir (CAC) supera 5x al de retener.",
        "Pérdida de clientes hacia la competencia (Altice).",
        "Viralidad: Las quejas públicas acumulan 4x likes.",
        "Oportunidad: Detección preventiva de averías."
    ], COLOR_GREEN)
    slide2.notes_slide.notes_text_frame.text = (
        "[00:35 - 01:15] ORLANDO: En un mercado de más de 5 millones de usuarios, la calidad percibida determina el churn. "
        "Las encuestas tradicionales tardan semanas, mientras en YouTube miles comentan caídas de fibra y demoras..."
    )

    # --- SLIDE 3: DATOS Y YOUTUBE DATA API V3 ---
    slide3 = crear_diapositiva_base(
        prs, "Ingesta de Datos: YouTube Data API v3",
        "Extracción ética, prudente de cuota y 100% reproducible",
        orador="Audric Rosario", apoyo_visual="Transición a VS Code: src/extractor_youtube.py"
    )
    agregar_tarjeta(slide3, 0.8, 1.8, 5.6, 4.8, "Métricas de Adquisición", [
        "Fuente Oficial: YouTube Data API v3 (Google Cloud).",
        "799 comentarios únicos reales recopilados.",
        "56 videos analizados (Institucionales, 5G, Fibra, Soporte).",
        "Consumo de Cuota: 649 unidades de 10,000 (Solo 6.49%).",
        "Pausas de cortesía y sanitización de caracteres."
    ], COLOR_CLARO_RED)
    agregar_tarjeta(slide3, 6.8, 1.8, 5.6, 4.8, "Reproducibilidad Offline", [
        "Datos crudos congelados en data/raw/ (CSV y JSON).",
        "Permite evaluación académica sin credenciales activas.",
        "Datos de prueba y mocks aislados en data/mocks/.",
        "Lectura segura de credenciales mediante .env (YT_API_KEY).",
        "Pipeline verificado bajo Python 3.12 y .venv local."
    ], COLOR_GREEN)
    slide3.notes_slide.notes_text_frame.text = (
        "[01:15 - 02:30] AUDRIC: Para este reto diseñamos un pipeline reproducible con herramientas 100% gratuitas. "
        "Consumimos solo 649 unidades de cuota y extrajimos 799 comentarios de 56 videos..."
    )

    # --- SLIDE 4: NLP Y TRANSFORMER EN ESPAÑOL ---
    slide4 = crear_diapositiva_base(
        prs, "Modelado NLP: Arquitectura Transformer en Español",
        "Clasificación contextual de sentimiento y categorización temática",
        orador="Audric Rosario", apoyo_visual="Transición a Jupyter Notebook / Celdas NLP"
    )
    agregar_tarjeta(slide4, 0.8, 1.8, 5.6, 4.8, "Preprocesamiento Lingüístico", [
        "Normalización fonética Unicode NFKD.",
        "Remoción de URLs, menciones @ y caracteres ruidosos.",
        "Tratamiento de jerga dominicana (nítido, avería, megas).",
        "Stopwords del dominio telco (claro, video, youtube)."
    ], COLOR_BORDER)
    agregar_tarjeta(slide4, 6.8, 1.8, 5.6, 4.8, "Modelo Transformer Preentrenado", [
        "Mecanismos de Autoatención Bidireccional (Self-Attention).",
        "Captura negaciones complejas, quejas técnicas y sarcasmo.",
        "Clasificación en 3 estados: Positivo, Negativo, Neutro.",
        "Cálculo de confianza y Net Sentiment Score (NSS).",
        "Segmentación en 5 áreas de servicio técnico."
    ], COLOR_CLARO_RED)
    slide4.notes_slide.notes_text_frame.text = (
        "[02:30 - 04:00] AUDRIC: El dialecto en redes en RD es complejo. Diseñamos un preprocesamiento fonético "
        "y aplicamos un Transformer en español con autoatención para capturar el contexto real..."
    )

    # --- SLIDE 5: HALLAZGOS Y DASHBOARD EJECUTIVO ---
    slide5 = crear_diapositiva_base(
        prs, "Hallazgos de Negocio y Dashboard Plotly",
        "Diagnóstico de reputación y focos críticos de fricción",
        orador="Orlando Benítez", apoyo_visual="Compartir Pantalla: dashboard_ejecutivo_claro.html interactivo"
    )
    agregar_tarjeta(slide5, 0.8, 1.8, 3.7, 4.8, "Indicadores Clave (KPIs)", [
        "Volumen Total: 799 opiniones auditadas.",
        "NSS Global: +5.76% (Favorable moderado).",
        "Positivos: 12.4% | Neutros: 81.0%.",
        "Negativos: 6.6% (Altamente concentrados)."
    ], COLOR_GREEN)
    agregar_tarjeta(slide5, 4.8, 1.8, 3.7, 4.8, "Liderazgo vs. Alerta", [
        "🟢 Red Móvil y 5G: +14.0% NSS (Gran admiración).",
        "🟢 Planes y Promos: +9.6% NSS.",
        "🔴 Atención al Cliente: -3.2% NSS (Crítico).",
        "🟡 Fibra Óptica: +4.4% NSS (Quejas nocturnas)."
    ], COLOR_CLARO_RED)
    agregar_tarjeta(slide5, 8.8, 1.8, 3.7, 4.8, "Patrón de Resonancia", [
        "Las quejas reciben 4x más likes que los elogios.",
        "Comunidad respalda reclamos por averías del 107.",
        "Términos críticos: lento, espera, ping, caída, bot.",
        "Riesgo directo de portabilidad hacia Altice."
    ], COLOR_BORDER)
    slide5.notes_slide.notes_text_frame.text = (
        "[04:00 - 05:30] ORLANDO: Aquí mostramos nuestro dashboard ejecutivo en Plotly. "
        "La red 5G es la fortaleza de Claro con +14% NSS, pero Atención al Cliente es la alerta roja con -3.2%..."
    )

    # --- SLIDE 6: SOLUCIÓN CLARO SENTINEL ---
    slide6 = crear_diapositiva_base(
        prs, "Solución Propuesta: 'Claro Sentinel NLP'",
        "De la escucha social pasiva al enrutamiento proactivo de quejas",
        orador="Orlando Benítez", apoyo_visual="Mostrar Diapositiva 6: Arquitectura Sentinel"
    )
    agregar_tarjeta(slide6, 0.8, 1.8, 3.7, 4.8, "1. Ingesta Continua", [
        "Worker automático consulta YouTube API.",
        "Detección de comentarios en videos oficiales y tutoriales.",
        "Filtro de privacidad y anonimización de datos."
    ], COLOR_BORDER)
    agregar_tarjeta(slide6, 4.8, 1.8, 3.7, 4.8, "2. Triaje Inteligente", [
        "Scoring de sentimiento con Transformer.",
        "Detección de averías técnicas críticas (ping, fibra caída).",
        "Priorización según intensidad y likes."
    ], COLOR_CLARO_RED)
    agregar_tarjeta(slide6, 8.8, 1.8, 3.7, 4.8, "3. Enlace al CRM", [
        "Generación de pre-ticket en Salesforce/Genesys.",
        "Respuesta pública oficial en < 2 horas.",
        "Reducción proyectada de churn del 1.8%."
    ], COLOR_GREEN)
    slide6.notes_slide.notes_text_frame.text = (
        "[05:30 - 06:05] ORLANDO: No basta con ver gráficos, se necesita actuar. "
        "Claro Sentinel detecta reclamos críticos y crea pre-tickets al CRM para responder en menos de 2 horas..."
    )

    # --- SLIDE 7: PRESUPUESTO Y GANTT ---
    slide7 = crear_diapositiva_base(
        prs, "Presupuesto de Implementación y Diagrama de Gantt",
        "Inversión cloud accesible y cronograma estructurado en 16 semanas",
        orador="Orlando Benítez", apoyo_visual="Mostrar Diapositiva 7: Tabla Presupuesto y Gantt"
    )
    agregar_tarjeta(slide7, 0.8, 1.8, 5.6, 4.8, "Presupuesto Anual (USD)", [
        "Servidor Inferencia AWS EC2 (g4dn.xlarge GPU): $2,160",
        "Base de Datos Cloud PostgreSQL (AWS RDS): $540",
        "YouTube Data API (GCP Free Tier): $0 (Gratuito)",
        "Capacitación al Personal de Atención Digital: $1,000",
        "Mantenimiento y MLOps Semestral: $1,150",
        "TOTAL ANUAL: $4,850 USD",
        "ROI: Se amortiza reteniendo 15 clientes/mes (< 90 días)."
    ], COLOR_GREEN)
    agregar_tarjeta(slide7, 6.8, 1.8, 5.6, 4.8, "Cronograma Gantt (16 Semanas)", [
        "Semanas 1-3: Requisitos y Conexión de Datos API.",
        "Semanas 4-7: Fine-tuning del Modelo Transformer.",
        "Semanas 8-10: Construcción de Dashboards en Looker/Plotly.",
        "Semanas 11-14: Piloto con el Equipo de CX y Soporte.",
        "Semanas 15-16: Despliegue General y Capacitación.",
        "Hito Clave: Salida a Producción (Go-Live) en Semana 16."
    ], COLOR_CLARO_RED)
    slide7.notes_slide.notes_text_frame.text = (
        "[06:05 - 06:45] ORLANDO: Implementar esta solución cuesta solo $4,850 dólares al año. "
        "Con retener 15 clientes al mes se amortiza en 90 días, con un plan secuencial de 16 semanas..."
    )

    # --- SLIDE 8: CONCLUSIONES Y CIERRE ---
    slide8 = crear_diapositiva_base(
        prs, "Conclusiones y Recomendaciones Estratégicas",
        "Valor tangible de la analítica aplicada a decisiones de telecomunicaciones",
        orador="Audric & Orlando", apoyo_visual="Cámaras de ambos oradores en pantalla dividida"
    )
    agregar_tarjeta(slide8, 0.8, 1.8, 5.6, 4.8, "Conclusiones Clave", [
        "1. Big Data Accesible: YouTube API + Transformers permite auditar percepción sin costos de licencia.",
        "2. Fortaleza Reputacional: La red 5G de Claro es líder en satisfacción nacional (+14.0% NSS).",
        "3. Brecha Operativa: Atención al cliente (-3.2% NSS) requiere intervención inmediata.",
        "4. Resonancia Negativa: El descontento se viraliza 4 veces más rápido que la satisfacción."
    ], COLOR_CLARO_RED)
    agregar_tarjeta(slide8, 6.8, 1.8, 5.6, 4.8, "Recomendaciones Gerenciales", [
        "1. Humanizar el Bot: Habilitar transferencia a humano en < 2 min ante quejas de averías.",
        "2. Estabilización de Fibra Nocturna: Auditar tráfico residencial en horas pico (7 a 11 PM).",
        "3. Monitoreo Competitivo: Extender el pipeline al canal de Altice RD para benchmarking.",
        "4. Enlace al CRM: Implementar Claro Sentinel para respuesta ágil."
    ], COLOR_GREEN)
    slide8.notes_slide.notes_text_frame.text = (
        "[06:45 - 07:30] AUDRIC: En conclusión, demostramos con rigor que herramientas de código abierto generan valor empresarial. "
        "ORLANDO: Recomendamos humanizar el bot y auditar la fibra. ¡Muchas gracias por su atención!"
    )

    prs.save(output_path)
    print(f"[EXITO] Presentación del Video Oficial generada en: {output_path}")
    return output_path


# ==============================================================================
# 2. PRESENTACIÓN RÁPIDA PARA LA CLASE (5 MINUTOS - 6 DIAPOSITIVAS)
# ==============================================================================
def generar_presentacion_clase_5min(output_path: str = "PRESENTACION_CLASE_5MIN.pptx"):
    print("[PPTX] Generando presentación para la Clase de 5 Minutos...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- SLIDE 1: PORTADA Y PITCH INICIAL (0:00 - 0:45) ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_DARK_NAVY
    bg1.line.color.rgb = COLOR_DARK_NAVY

    top1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.15))
    top1.fill.solid()
    top1.fill.fore_color.rgb = COLOR_CLARO_RED

    tx = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
    tf = tx.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "UAPA • APLICACIONES ANALÍTICAS DE BIG DATA • PROYECTO FINAL"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MUTED

    p2 = tf.add_paragraph()
    p2.text = "Auditoría de Experiencia y Sentimiento de Marca en Claro República Dominicana"
    p2.font.size = Pt(30)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.space_before = Pt(16)

    p3 = tf.add_paragraph()
    p3.text = "Transformando la Voz de Redes Sociales (YouTube API) en Retención y Calidad de Servicio mediante NLP"
    p3.font.size = Pt(15)
    p3.font.color.rgb = COLOR_CLARO_RED
    p3.space_before = Pt(8)

    p4 = tf.add_paragraph()
    p4.text = "Facilitador: Luis Eduardo Bayonet Robles\nExpositores: Audric André Rosario Rosario (100089140) & Orlando Benítez Ventura (100090873) | Tiempo: 5 Minutos"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_WHITE
    p4.space_before = Pt(22)

    slide1.notes_slide.notes_text_frame.text = (
        "[0:00 - 0:45] ORLANDO: Buenos días profesor y compañeros. Hoy les presentamos cómo la analítica de redes sociales "
        "y el Big Data salvan millones de pesos en retención de clientes para Claro República Dominicana..."
    )

    # --- SLIDE 2: EL RETO EMPRESARIAL (0:45 - 1:30) ---
    slide2 = crear_diapositiva_base(
        prs, "El Reto: Puntos Ciegos en la Satisfacción de Clientes",
        "Por qué las encuestas tradicionales ya no son suficientes",
        orador="Orlando Benítez (0:45 - 1:30)", apoyo_visual="Diapositiva 2"
    )
    agregar_tarjeta(slide2, 0.8, 1.8, 5.6, 4.8, "La Falla de las Encuestas Tradicionales", [
        "El 96% de los usuarios NO responde encuestas de satisfacción.",
        "Reportes mensuales tardíos: cuando llegan, el cliente ya se cambió.",
        "Costo elevado y muestreo desconectado de la operación diaria.",
        "Falta de visibilidad sobre quejas reales de averías de internet."
    ], COLOR_BORDER)
    agregar_tarjeta(slide2, 6.8, 1.8, 5.6, 4.8, "La Mina de Oro de YouTube", [
        "Miles de comentarios voluntarios en anuncios, tutoriales y comparativas.",
        "La voz honesta y sin filtro del usuario dominicano.",
        "Información en tiempo real sobre fallas de fibra y red 5G.",
        "Problema: Los datos viven dispersos y sin análisis estructurado."
    ], COLOR_CLARO_RED)
    slide2.notes_slide.notes_text_frame.text = (
        "[0:45 - 1:30] ORLANDO: Claro gasta fortunas en encuestas que el 96% no responde. "
        "Sin embargo, en videos de YouTube los usuarios dejan miles de opiniones sobre fibra y 5G. Cerramos esa brecha..."
    )

    # --- SLIDE 3: ARQUITECTURA TÉCNICA Y DATOS (1:30 - 2:30) ---
    slide3 = crear_diapositiva_base(
        prs, "Arquitectura Big Data: Extracción y NLP Transformer",
        "Herramientas 100% gratuitas con consumo prudente de API",
        orador="Audric Rosario (1:30 - 2:30)", apoyo_visual="Diapositiva 3"
    )
    agregar_tarjeta(slide3, 0.8, 1.8, 5.6, 4.8, "Extracción YouTube Data API v3", [
        "799 comentarios únicos reales recolectados.",
        "56 videos corporativos y comparativas técnicas.",
        "Prudencia estricta: Solo 649 unidades de 10,000 diarias (< 7%).",
        "Datos locales congelados en data/raw/ para reproducibilidad."
    ], COLOR_GREEN)
    agregar_tarjeta(slide3, 6.8, 1.8, 5.6, 4.8, "Modelo Transformer en Español", [
        "Preprocesamiento fonético adaptado a dialecto dominicano.",
        "Autoatención bidireccional para capturar ironía y contexto.",
        "Clasificación multiclase: Positivo, Negativo, Neutro.",
        "Mapeo a 5 tópicos: Fibra, 5G, Soporte, Facturas, Planes."
    ], COLOR_CLARO_RED)
    slide3.notes_slide.notes_text_frame.text = (
        "[1:30 - 2:30] AUDRIC: Utilizamos herramientas 100% gratuitas. Conectamos la YouTube Data API v3 "
        "y extrajimos 799 comentarios reales de 56 videos. Construimos un pipeline NLP con Transformers en español..."
    )

    # --- SLIDE 4: HALLAZGOS Y DASHBOARD EJECUTIVO (2:30 - 3:30) ---
    slide4 = crear_diapositiva_base(
        prs, "Hallazgos Clave: Diagnóstico y Riesgo de Viralidad",
        "Métricas del Dashboard Ejecutivo interactivo en Plotly",
        orador="Audric Rosario (2:30 - 3:30)", apoyo_visual="Diapositiva 4"
    )
    agregar_tarjeta(slide4, 0.8, 1.8, 5.6, 4.8, "Distribución y Net Sentiment Score", [
        "Net Sentiment Score (NSS): +5.76% (Percepción favorable).",
        "Neutros: 81.0% (Consultas de cobertura, precios y ayuda).",
        "Positivos: 12.4% (Orgullo por velocidad 5G y cobertura nacional).",
        "Negativos: 6.6% (Demoras en call center y fibra nocturna)."
    ], COLOR_GREEN)
    agregar_tarjeta(slide4, 6.8, 1.8, 5.6, 4.8, "Foco Crítico y Resonancia", [
        "🔴 Atención al Cliente es la mayor alarma: -3.2% NSS.",
        "🟢 Red 5G es la mayor fortaleza: +14.0% NSS.",
        "⚠️ Las quejas reciben 4 VECES MÁS LIKES que los elogios.",
        "El descontento se viraliza y acelera el churn hacia Altice."
    ], COLOR_CLARO_RED)
    slide4.notes_slide.notes_text_frame.text = (
        "[2:30 - 3:30] AUDRIC: El NSS es de +5.8%. El 81% son consultas de clientes. La red 5G es la joya con +14%, "
        "pero Atención al Cliente está en rojo con -3.2%. Además las quejas reciben 4 veces más apoyo comunitario..."
    )

    # --- SLIDE 5: SOLUCIÓN CLARO SENTINEL Y ROI (3:30 - 4:30) ---
    slide5 = crear_diapositiva_base(
        prs, "Solución 'Claro Sentinel NLP' y Retorno de Inversión",
        "Resolución proactiva en menos de 2 horas vinculada al CRM",
        orador="Orlando Benítez (3:30 - 4:30)", apoyo_visual="Diapositiva 5"
    )
    agregar_tarjeta(slide5, 0.8, 1.8, 5.6, 4.8, "Funcionamiento del Sistema", [
        "Escucha continua en canales oficiales de YouTube.",
        "Clasificación automática de gravedad con Transformer.",
        "Enrutamiento directo a la mesa de ayuda (CRM Salesforce).",
        "Respuesta pública ágil en el mismo canal en < 2 horas."
    ], COLOR_CLARO_RED)
    agregar_tarjeta(slide5, 6.8, 1.8, 5.6, 4.8, "Presupuesto y Retorno (ROI)", [
        "Costo total anual: $4,850 USD en infraestructura cloud.",
        "Aprovecha el Free Tier de YouTube Data API.",
        "Amortización en < 90 días reteniendo solo 15 clientes/mes.",
        "ROI superior al 300% en el primer año de operación."
    ], COLOR_GREEN)
    slide5.notes_slide.notes_text_frame.text = (
        "[3:30 - 4:30] ORLANDO: Diseñamos Claro Sentinel para enlazar quejas al CRM en menos de 2 horas. "
        "El costo es de solo $4,850 dólares al año. Con retener 15 clientes al mes se paga en 90 días..."
    )

    # --- SLIDE 6: CONCLUSIONES Y CIERRE (4:30 - 5:00) ---
    slide6 = crear_diapositiva_base(
        prs, "Conclusiones, Recomendaciones y Preguntas",
        "Impacto estratégico para Claro República Dominicana",
        orador="Audric & Orlando (4:30 - 5:00)", apoyo_visual="Diapositiva 6"
    )
    agregar_tarjeta(slide6, 0.8, 1.8, 5.6, 4.8, "Conclusiones del Proyecto", [
        "Herramientas gratuitas permiten analítica de calidad corporativa.",
        "La red 5G es el gran diferenciador positivo de Claro RD.",
        "La atención al cliente virtual requiere humanización urgente.",
        "Código y datasets disponibles y reproducibles en GitHub."
    ], COLOR_BORDER)
    agregar_tarjeta(slide6, 6.8, 1.8, 5.6, 4.8, "Recomendaciones Gerenciales", [
        "Reducir los menús del bot a máximo 2 minutos para transferir a agente.",
        "Auditar y estabilizar la fibra óptica residencial en horario de 7 a 11 PM.",
        "Monitorear la competencia (Altice RD) en tiempo real.",
        "¡Quedamos a su disposición para la sesión de preguntas!"
    ], COLOR_GREEN)
    slide6.notes_slide.notes_text_frame.text = (
        "[4:30 - 5:00] AUDRIC: Demostramos que con herramientas gratuitas se construye analítica de clase mundial. "
        "ORLANDO: Muchas gracias por su atención, abrimos la sesión de preguntas."
    )

    prs.save(output_path)
    print(f"[EXITO] Presentación de la Clase de 5 Minutos generada en: {output_path}")
    return output_path


if __name__ == "__main__":
    generar_presentacion_video()
    generar_presentacion_clase_5min()
